from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation from template
import os
template_path = r"C:\Users\kusha\Downloads\[EXT] Solution Challenge 2026 - Prototype PPT Template.pptx"
prs = Presentation(template_path)

# Helper function to clear and set text
def set_text(ph, text, font_size=18, bold=False, color=None):
    ph.text = text
    for paragraph in ph.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = RGBColor(*color)

# SLIDE 2 - Team Details
slide = prs.slides[1]
for shape in slide.shapes:
    if shape.has_text_frame:
        if "Team name:" in shape.text:
            set_text(shape.text_frame, "Team name: RAI Innovators", 18, True)
        elif "Team leader name:" in shape.text:
            set_text(shape.text_frame, "Team leader name: Kushal Chala", 18, True)
        elif "Problem Statement:" in shape.text:
            set_text(shape.text_frame, """Problem Statement:
AI systems deployed in healthcare, hiring, lending, and insurance 
exhibit demographic bias that goes undetected due to lack of 
accessible open-source auditing tools for organizations without 
specialized ML expertise.""", 16)

# SLIDE 3 - Brief about solution
slide = prs.slides[2]
for shape in slide.shapes:
    if shape.has_text_frame and "Brief" in shape.text:
        set_text(shape.text_frame, """RAI (Responsible AI Audit Platform) is an open-source, web-based fairness auditing tool that enables organizations to detect, measure, and correct bias in their datasets and AI models.

Users upload datasets from 4 key domains (Healthcare, Hiring, Lending, Insurance), and the platform automatically:
• Detects demographic bias using Demographic Parity metrics
• Visualizes group-level positive rates with interactive charts
• Applies mitigation via reweighting algorithms
• Generates compliance-ready audit reports

Built with FastAPI backend + React frontend, deployed on Render (backend) and Vercel (frontend), with sample datasets matching 2024 US industry distributions (EEOC, FFIEC, CMS).""", 16)

# SLIDE 4 - Opportunities
slide = prs.slides[3]
for shape in slide.shapes:
    if shape.has_text_frame:
        if "different" in shape.text.lower():
            set_text(shape.text_frame, """How different from existing ideas?
• Most bias auditing tools (AIF360, Fairlearn) require Python expertise
• Commercial tools (Fiddler, Arthur) are closed-source and expensive
• RAI is zero-setup: browser-based, pre-loaded with 4 domain datasets
• Unique feature: One-click sample dataset loading for instant demos""", 16)
        elif "solve the problem" in shape.text.lower():
            set_text(shape.text_frame, """How will it solve the problem?
• Auto-detects sensitive attributes (gender, race) and outcome columns
• Applies 80% four-fifths rule (EEOC standard) automatically
• Quantifies bias with Demographic Parity Ratio & Difference metrics
• Provides actionable mitigation with before/after comparison charts""", 16)
        elif "USP" in shape.text:
            set_text(shape.text_frame, """USP of the proposed solution:
• Open-source + zero-code interface
• Pre-built datasets for 4 high-impact domains
• Built-in regulatory compliance (EU AI Act, EEOC, India DPDPA)
• Realtime fairness visualization with group-level positive rates
• AI-powered audit summaries using Google Gemini""", 16)

# SLIDE 5 - List of features
slide = prs.slides[4]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """1. Dataset Upload — Drag & drop CSV/Excel files (10MB limit)
2. Sample Datasets — 4 pre-loaded domain datasets (2500 records each)
   • Healthcare (treatment approval bias)
   • Hiring (recruitment bias)
   • Lending (loan approval bias)
   • Insurance (policy approval bias)
3. Auto-Detection — Sensitive attributes & target columns auto-identified
4. Fairness Metrics — Demographic Parity Ratio & Difference (80% threshold)
5. Interactive Charts — Group-level positive rate visualization
6. Bias Mitigation — Reweighting algorithm with before/after comparison
7. Compliance Reports — Downloadable audit reports (HTML)
8. Mitigated Dataset Export — Download bias-corrected CSV
9. AI-Powered Summary — Google Gemini generates audit insights""", 16)

# SLIDE 6 - Process flow diagram
slide = prs.slides[5]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """[User] → [Upload Dataset / Load Sample]
         ↓
    [Auto-Detect Columns]
         ↓
    [Select Sensitive Attrs + Target]
         ↓
    [Run Fairness Analysis]
         ↓
    [View Results: DPR, DPD, Group Charts]
         ↓
    [Optional: Apply Mitigation]
         ↓
    [Download Report / Mitigated CSV]

Use Cases:
• HR teams auditing recruitment data for gender/race bias
• Banks checking loan approval fairness
• Hospitals validating treatment approval processes
• Insurance companies auditing policy decisions""", 14)

# SLIDE 7 - Wireframes (optional)
slide = prs.slides[6]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """Screenshots to include:
1. Hero Section — Animated crowd scene with "RAI" title, domain pills
2. Upload Section — Drag & drop zone + sample dataset cards
3. Column Selector — Multi-select sensitive attrs, target variable
4. Results Panel — DPR/DPD metric cards + group rate bar chart
5. Mitigation Section — Before/after comparison + download buttons
6. AI Summary — Gemini-powered audit insights (auto-generated)""", 16)

# SLIDE 8 - Architecture diagram
slide = prs.slides[7]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """┌─────────────────┐     HTTPS     ┌──────────────────────┐
│  React Frontend │ ←─────────→ │  FastAPI Backend     │
│  (Vercel)      │             │  (Render)            │
│                 │             │                      │
│ • UploadSection │             │ • /upload           │
│ • ColumnSelector│             │ • /analyze          │
│ • MetricsPanel  │             │ • /mitigate         │
│ • MitigationSec │             │ • /sample-datasets  │
└─────────────────┘             │ • /gemini-summary   │
                                │ • /download/mitigated│
                                └──────────────────────┘
                                          ↓
                                ┌──────────────────────┐
                                │  In-Memory Storage   │
                                │  (DATASET_STORAGE)  │
                                │  (AUDIT_STORAGE)    │
                                └──────────────────────┘

Tech Stack:
Frontend: React 18, Vite, Framer Motion, Recharts, react-dropzone
Backend: FastAPI, Pandas, NumPy, Google Gemini API
Deployment: Vercel (frontend), Render (backend)""", 14)

# SLIDE 9 - Technologies used
slide = prs.slides[8]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """Frontend:
• React 18 — UI framework
• Vite — Build tool
• Framer Motion — Animations
• Recharts — Data visualization
• Axios — API client
• react-dropzone — File upload

Backend:
• FastAPI — REST API framework
• Pandas — Dataset processing
• NumPy — Numerical computations
• Google Generative AI — Gemini integration
• Uvicorn — ASGI server

Deployment & DevOps:
• Vercel — Frontend hosting
• Render — Backend hosting
• Git/GitHub — Version control""", 16)

# SLIDE 10 - Estimated implementation cost
slide = prs.slides[9]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """Development Time: ~40 hours

Technologies Used: All open-source (no licensing costs)

Hosting Costs:
• Render (backend): Free tier (or ~$7/month for always-on)
• Vercel (frontend): Free tier
• Total Monthly Cost: $0–$7 USD

Optional Scaling:
• Render Starter: $7/month
• Custom domain: ~$12/year

Google Gemini API: Free tier (15 RPM, 1M tokens/min)""", 16)

# SLIDE 11 - Snapshots of MVP
slide = prs.slides[10]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """Include 4–5 screenshots:
1. Hero section with animated crowd + domain pills (Healthcare, Hiring, etc.)
2. Upload area with sample dataset cards (4 domain options)
3. Column selector with detected sensitive attributes
4. Results panel showing DPR = 72.3% (FAIL) with bar chart
5. Mitigation section with before/after comparison chart
6. Gemini AI summary with actionable recommendations""", 16)

# SLIDE 12 - Additional Details/Future Development
slide = prs.slides[11]
for shape in slide.shapes:
    if shape.has_text_frame:
        set_text(shape.text_frame, """Current Capabilities:
• CSV and Excel file support
• 4 industry-domain sample datasets (2500 records each)
• Demographic Parity metrics (DPR, DPD)
• Reweighting-based bias mitigation
• Compliance with EEOC 4/5ths rule, EU AI Act
• Google Gemini-powered audit summaries

Future Development:
• Enhanced Gemini integration for real-time audit recommendations
• Additional fairness metrics (Equal Opportunity, Calders & Verwer)
• Support for multi-class and regression outcomes
• PostgreSQL database for persistent audit history
• PDF report generation with executive summaries
• Role-based access control for enterprise teams
• API endpoints for CI/CD pipeline integration""", 16)

# SLIDE 13 - Links
slide = prs.slides[12]
for shape in slide.shapes:
    if shape.has_text_frame:
        if "GitHub" in shape.text:
            set_text(shape.text_frame, """GitHub Public Repository:
https://github.com/kushalchalla981-tech/rai-responsible.ai

Demo Video Link (3 Minutes):
[Upload to YouTube and paste link here]

MVP Link:
https://rai-responsible-77cdcaaop-kushalchalla981-techs-projects.vercel.app/

Working Prototype Link:
Frontend: https://rai-responsible-77cdcaaop-kushalchalla981-techs-projects.vercel.app/
Backend API: https://rai-responsible-ai.onrender.com
API Docs: https://rai-responsible-ai.onrender.com/docs""", 14)

# Save the presentation
output_path = r"C:\projectwa\new\rai-audit-platform\RAI_Audit_Platform_Submission.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
